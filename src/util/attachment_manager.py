import os
import zlib
import uuid
import base64

from dotenv import load_dotenv
from flask import Flask, request, jsonify, send_from_directory, Response, stream_with_context
from flask_cors import CORS
import fitz  # PyMuPDF
from docx import Document
import olefile
import csv
from pptx import Presentation
from openpyxl import load_workbook
from flask import send_from_directory

from config.settings import *
from util.file_manager import _sanitize_filename

load_dotenv("src/parquet/.env")

# PDF 파일에서 텍스트 추출
def _extract_text_from_pdf(file_path):
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        print(f"[PDF Extract Error] {e}")
    return text

# Word 파일에서 텍스트 추출
def _extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"[Docx Extract Error] {e}")
    return text

# HWP 파일에서 텍스트 추출
def _extract_text_from_hwp(file_path):
    text = ""
    try:
        f = olefile.OleFileIO(file_path)
        dirs = f.listdir()
        sections = [d for d in dirs if "BodyText/Section" in "/".join(d)]
        for section in sections:
            stream = f.openstream("/".join(section))
            data = stream.read()
            try:
                decompressed = zlib.decompress(data, -15)
                decoded_text = decompressed.decode("utf-16-le", errors="ignore")
                clean_text = "".join(c for c in decoded_text if c.isalnum() or c in " \n\t.,()[]")
                text += clean_text + "\n"
            except Exception as e:
                print(f"[HWP Decode Error in {section}] {e}")
        f.close()
    except Exception as e:
        print(f"[HWP Extract Error] {e}")
    return text

# TXT 파일에서 텍스트 추출
def _extract_text_from_txt(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="cp949") as f:
                text = f.read()
        except Exception as e:
            print(f"[TXT Extract Error] {e}")
    except Exception as e:
        print(f"[TXT Extract Error] {e}")
    return text

# PPTX 파일에서 텍스트 추출
def _extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[PPTX Extract Error] {e}")
    return text

# XLSX 파일에서 텍스트 추출
def _extract_text_from_xlsx(file_path):
    text = ""
    try:
        wb = load_workbook(file_path, data_only=True)
        for ws in wb.worksheets:
            text += f"[Sheet] {ws.title}\n"
            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    text += " | ".join(row_values) + "\n"
            text += "\n"
    except Exception as e:
        print(f"[XLSX Extract Error] {e}")
    return text

# CSV 파일에서 텍스트 추출
def _extract_text_from_csv(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                row_values = [str(cell) if cell is not None else "" for cell in row]
                text += " | ".join(row_values) + "\n"
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="cp949", newline="") as f:
                reader = csv.reader(f)
                for row in reader:
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    text += " | ".join(row_values) + "\n"
        except Exception as e:
            print(f"[CSV Extract Error] {e}")
    except Exception as e:
        print(f"[CSV Extract Error] {e}")
    return text

# attachment payload에서 base64를 받아 서버 로컬에 파일 저장
def _save_attachment_from_base64(file_info: dict, save_dir: str) -> tuple[str, str]:
    original_name = file_info.get("name") or "attachment.bin"
    safe_name = _sanitize_filename(original_name)
    mail_id = str(file_info.get("mail_id") or "no_mail_id")
    data_base64 = file_info.get("data_base64") or ""

    if not data_base64:
        raise ValueError(f"attachment data_base64 missing: {original_name}")

    os.makedirs(save_dir, exist_ok=True)

    ext = os.path.splitext(safe_name)[1].lower()
    unique_name = f"{mail_id}_{uuid.uuid4().hex[:8]}{ext or '.bin'}"
    saved_path = os.path.join(save_dir, unique_name)

    if "," in data_base64 and "base64" in data_base64[:100]:
        data_base64 = data_base64.split(",", 1)[1]

    file_bytes = base64.b64decode(data_base64)

    with open(saved_path, "wb") as f:
        f.write(file_bytes)

    return saved_path, original_name

